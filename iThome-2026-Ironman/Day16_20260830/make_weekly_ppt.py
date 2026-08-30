# -*- coding: utf-8 -*-
"""把 Stock_Summary/ 的五份每日財經節目摘要，整理成一份「上週回顧與下週判斷」簡報。

內容規則沿用 Day 3-4 那份 YouTube 摘要規則書：
  - 只寫來源明確講過的內容，不推論、不用自己的財經知識補空缺。
  - 只轉述，不評論；全簡報不得出現我自己的判斷、建議或警語。
  - 每一條都要能回溯到某一位分析師，寫出姓名。
  - 態度只用五個詞：看多／偏多／中性／偏空／看空。
  - 來源沒查證到的股名或代號一律標「待確認」，不憑印象填。
差別只在輸出載體：規則書輸出的是可以貼到 LINE 的純文字，這裡輸出的是 pptx。

視覺系統與 design/ 底下那組設計稿一致：深藍 #0C2340 為主，台股慣例紅漲綠跌，
數字一律走等寬字（Consolas）好對齊，每頁一條頁碼列與一行來源註記。

用法（python-pptx 只裝在 office 這個環境）：
    C:/Users/Alex/anaconda3/envs/office/python.exe make_weekly_ppt.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
OUT_FILE = HERE / "Stock_Weekly_Review_20260824-0828.pptx"

# ── 色票（與 design/ 的設計稿同一組）──────────────────────────
NAVY = RGBColor(0x0C, 0x23, 0x40)
NAVY_SOFT = RGBColor(0x9F, 0xB3, 0xCC)
NAVY_TEXT = RGBColor(0xC7, 0xD5, 0xE5)
PAPER = RGBColor(0xF6, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x10, 0x18, 0x28)
MUTED = RGBColor(0x5B, 0x64, 0x72)
LINE = RGBColor(0xD8, 0xDE, 0xE7)
UP = RGBColor(0xC1, 0x12, 0x1F)        # 台股：紅為漲
UP_SOFT = RGBColor(0xE9, 0xB7, 0xBB)
DOWN = RGBColor(0x0F, 0x7A, 0x55)      # 台股：綠為跌
PINK = RGBColor(0xFF, 0xD9, 0xDC)

SANS = "Microsoft JhengHei"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.58)
BODY_W = SLIDE_W - MARGIN * 2
BAND_H = Inches(1.02)


# ── 基本元件 ────────────────────────────────────────────────

def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def blank(prs, *, cover=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY if cover else PAPER)
    return slide


def textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return frame


def write(frame, text, *, size=15, color=INK, bold=False, space_after=5,
          first=False, align=PP_ALIGN.LEFT, mono=False, space_before=0):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.alignment = align
    para.space_after = Pt(space_after)
    para.space_before = Pt(space_before)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = MONO if mono else SANS
    return para


def band(slide, number, title, subtitle=None):
    """每頁頂端的深藍標題列。"""
    rect(slide, 0, 0, SLIDE_W, BAND_H, fill=NAVY)
    frame = textbox(slide, MARGIN, Inches(0.26), Inches(1.0), Inches(0.4))
    write(frame, number, size=14, color=UP, bold=True, mono=True, first=True, space_after=0)
    frame = textbox(slide, MARGIN + Inches(0.62), Inches(0.2), Inches(8.4), Inches(0.6))
    write(frame, title, size=26, color=WHITE, bold=True, first=True, space_after=0)
    if subtitle:
        frame = textbox(slide, Inches(8.6), Inches(0.34), SLIDE_W - Inches(8.6) - MARGIN,
                        Inches(0.4))
        write(frame, subtitle, size=12, color=NAVY_SOFT, first=True, space_after=0,
              align=PP_ALIGN.RIGHT)


def footer(slide, text):
    frame = textbox(slide, MARGIN, Inches(6.98), BODY_W, Inches(0.32))
    write(frame, text, size=9.5, color=MUTED, first=True, space_after=0)


def card(slide, left, top, width, height, *, accent=None, side="top"):
    """白底卡片，可在上緣或左緣加一條強調色。"""
    shape = rect(slide, left, top, width, height, fill=WHITE)
    if accent is not None:
        if side == "top":
            rect(slide, left, top, width, Pt(3), fill=accent)
        else:
            rect(slide, left, top, Pt(3), height, fill=accent)
    return shape


def table(slide, left, top, width, col_widths, rows, *, font=12.5, row_height=Inches(0.34),
          header_font=12):
    """深藍表頭 + 白／灰交錯列。"""
    shape = slide.shapes.add_table(len(rows), len(col_widths), left, top, width,
                                   row_height * len(rows))
    tbl = shape.table
    tbl.horz_banding = False
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Emu(int(w))
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(0.32) if r == 0 else row_height
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else PAPER
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = Pt(header_font if r == 0 else font)
                run.font.bold = r == 0
                run.font.name = SANS
                run.font.color.rgb = WHITE if r == 0 else INK
    return tbl


# ── 各頁 ────────────────────────────────────────────────────

def cover(prs):
    slide = blank(prs, cover=True)
    rect(slide, MARGIN, Inches(1.05), Inches(0.48), Pt(4), fill=UP)
    frame = textbox(slide, MARGIN + Inches(0.62), Inches(0.96), Inches(6), Inches(0.34))
    write(frame, "W E E K L Y   M A R K E T   R E V I E W", size=12, color=NAVY_SOFT,
          first=True, space_after=0)

    frame = textbox(slide, MARGIN, Inches(1.6), Inches(11), Inches(2.4))
    write(frame, "台股財經節目週摘要", size=50, color=WHITE, bold=True, first=True, space_after=10)
    write(frame, "2026 / 08 / 24  —  08 / 28", size=25, color=WHITE, bold=True, mono=True,
          space_after=6)
    write(frame, "上週回顧　×　下週判斷", size=19, color=NAVY_TEXT, bold=True, space_after=0)

    stats = [
        ("來源", "5", "個交易日的節目摘要檔"),
        ("涵蓋", "約 30", "位分析師的發言，逐日彙整後跨日比對"),
    ]
    for i, (label, value, note) in enumerate(stats):
        left = MARGIN + Inches(i * 3.1)
        frame = textbox(slide, left, Inches(4.55), Inches(2.9), Inches(1.3))
        write(frame, label, size=11.5, color=RGBColor(0x7E, 0x95, 0xAF), first=True, space_after=3)
        write(frame, value, size=30, color=WHITE, bold=True, mono=True, space_after=3)
        write(frame, note, size=12.5, color=NAVY_TEXT, space_after=0)

    rect(slide, Inches(6.9), Inches(4.55), Pt(1), Inches(1.25), fill=RGBColor(0x2B, 0x44, 0x63))
    frame = textbox(slide, Inches(7.25), Inches(4.55), Inches(5.4), Inches(1.4))
    write(frame, "整理原則", size=11.5, color=RGBColor(0x7E, 0x95, 0xAF), first=True, space_after=3)
    write(frame, "只轉述節目講過的內容，不加入整理者自己的判斷；每一條都標示是哪位分析師講的；"
                 "查不到的股名與代號標「待確認」。", size=13, color=RGBColor(0xE4, 0xEA, 0xF1),
          space_after=0)

    rect(slide, MARGIN, Inches(6.32), BODY_W, Pt(1), fill=RGBColor(0x2B, 0x44, 0x63))
    rect(slide, MARGIN, Inches(6.6), Inches(0.14), Inches(0.14), fill=UP)
    frame = textbox(slide, MARGIN + Inches(0.3), Inches(6.52), Inches(8), Inches(0.4))
    write(frame, "本簡報為節目內容整理，不構成任何投資建議。", size=16, color=PINK, bold=True,
          first=True, space_after=0)
    frame = textbox(slide, Inches(9.5), Inches(6.58), Inches(3.2), Inches(0.3))
    write(frame, "Generated with Claude Code", size=9.5, color=RGBColor(0x7E, 0x95, 0xAF),
          first=True, space_after=0, align=PP_ALIGN.RIGHT)


def week_track(prs):
    slide = blank(prs)
    band(slide, "01", "上週指數軌跡", "先看發生了什麼，再看誰講對了什麼")
    rows = [
        ["日期", "收盤／變動", "成交量", "當日主軸"],
        ["08/24 (一)", "44,762　-461 點", "6,294 億（今年 4 月以來新低量）",
         "極度觀望；矽晶圓全面漲價成為弱勢盤中亮點"],
        ["08/25 (二)", "破底後尾盤大拉近千點", "量縮觀望", "44,000 點支撐守住，等輝達財報與 PCE"],
        ["08/26 (三)", "+663 點", "放大至 8,000 億以上",
         "美債殖利率與油價大跌，資金回流電子；貨櫃航運重挫"],
        ["08/27 (四)", "45,975　+142 點（收長上影線）", "量能不足",
         "輝達財報超標，盤中漲逾 500 點挑戰前高後收斂"],
        ["08/28 (五)", "四天自低點累計 +2,300 點", "突破兆元",
         "外資 8 月買超 3,707 億元創單月新高，台幣升至 31.57"],
    ]
    widths = [Inches(1.45), Inches(2.55), Inches(2.6), Inches(5.57)]
    tbl = table(slide, MARGIN, Inches(1.42), BODY_W, widths, rows, row_height=Inches(0.56))
    # 漲跌數字上色：綠跌、紅漲
    for r, color in ((1, DOWN), (3, UP), (4, UP), (5, UP)):
        for run in tbl.cell(r, 1).text_frame.paragraphs[0].runs:
            run.font.color.rgb = color
            run.font.bold = True
            run.font.name = MONO
    for r in range(1, 6):
        for run in tbl.cell(r, 0).text_frame.paragraphs[0].runs:
            run.font.name = MONO
            run.font.bold = True

    box = rect(slide, MARGIN, Inches(4.92), BODY_W, Inches(1.32), fill=NAVY)
    frame = textbox(slide, MARGIN + Inches(0.28), Inches(5.1), BODY_W - Inches(0.56), Inches(1.0))
    write(frame, "週一恐慌新低量　→　週二尾盤反攻　→　週三資金回流電子　→　"
                 "週四輝達利多但收上影線　→　週五量能過兆但未過前高",
          size=16, color=WHITE, bold=True, first=True, space_after=6)
    write(frame, "陳昆仁在週一即指出「新低量不等於新低價」，是主力量縮洗浮額；"
                 "四天後大盤自低點漲逾 2,000 點。", size=13, color=NAVY_TEXT, space_after=0)
    footer(slide, "指數與量能數字均為節目中提及者，未經另行查證。")


def macro(prs):
    slide = blank(prs)
    band(slide, "02", "總體環境：三個變數同時轉向", "利率、油價、匯率")

    cards = [
        (DOWN, "美債殖利率", ["10 年期約 4.73%", "30 年期約 5.27%"], "↓　長天期殖利率下滑",
         "財政部擴大長債回購，規模約 1 兆美元；利率走跌有助高本益比成長股估值修復。",
         "游庭皓、林漢偉、容逸燊"),
        (DOWN, "油價與通膨", ["關稅戰擴大至加拿大", "通膨預期升溫"], "↓　布蘭特原油單日大跌約 4%",
         "伊朗與阿曼協商重啟霍爾木茲海峽臨時航道並排雷，全球通膨預期降溫。",
         "林漢偉、劉育綸、蔚辰"),
        (UP, "匯率與外資", ["外資大舉匯出 280 億美元"], "↑　台幣升至 31.57",
         "美元指數走弱，外資期現貨同步翻多，8 月買超 3,707 億元，創史上單月最高。",
         "容逸燊、高憲容"),
    ]
    card_w = Inches(4.02)
    for i, (accent, title, before, after, body, who) in enumerate(cards):
        left = MARGIN + Inches(i * 4.16)
        card(slide, left, Inches(1.42), card_w, Inches(2.9), accent=accent)
        frame = textbox(slide, left + Inches(0.24), Inches(1.62), card_w - Inches(0.48),
                        Inches(2.6))
        write(frame, title, size=18, color=INK, bold=True, first=True, space_after=8)
        write(frame, "週初", size=10.5, color=MUTED, space_after=3)
        for line in before:
            write(frame, line, size=14, color=INK, mono=("%" in line), space_after=2)
        write(frame, after, size=14, color=accent, bold=True, space_after=6, space_before=6)
        write(frame, body, size=12.5, color=INK, space_after=4)
        write(frame, who, size=11.5, color=MUTED, space_after=0)

    card(slide, MARGIN, Inches(4.55), BODY_W, Inches(1.72), accent=NAVY, side="left")
    frame = textbox(slide, MARGIN + Inches(0.28), Inches(4.74), BODY_W - Inches(0.56), Inches(0.3))
    write(frame, "其他總經數字", size=15, color=NAVY, bold=True, first=True, space_after=0)
    notes = [
        ("7 月 PCE 年增 3.7%、核心 3.3%，與 6 月持平，通膨具黏滯性；9 月大機率按兵不動",
         "胡睿涵、祥維、紀緯明"),
        ("台灣 7 月工業生產與製造業生產指數創單月歷史新高，景氣燈號連八紅", "蔚辰、游庭皓"),
        ("大摩上修台灣今年經濟成長率至 11.6%，為 40 年來最快", "胡睿涵"),
    ]
    for i, (text, who) in enumerate(notes):
        left = MARGIN + Inches(0.28) + Inches(i * 3.98)
        frame = textbox(slide, left, Inches(5.2), Inches(3.7), Inches(1.0))
        write(frame, text, size=12.5, color=INK, first=True, space_after=4)
        write(frame, who, size=11, color=MUTED, space_after=0)
    footer(slide, "以上為節目中提及的數據，未另行查證原始統計來源。")


def consensus_up(prs):
    slide = blank(prs)
    band(slide, "03", "一週共識：看多", "以「五天當中有幾天被列為一致看多」計算")
    rows = [
        ["產業／主題", "天數", "代表看多者"],
        ["CPO 光通訊／矽光子", "4", "林漢偉、陳昆仁、劉育綸、容逸燊、黃豐凱"],
        ["ABF 載板／CCL／PCB", "4", "林漢偉、劉育綸、陳於晨、陳威良、囿羽"],
        ["記憶體／HBM", "4", "林漢偉、陳昆仁、高憲容、老王"],
        ["先進封裝設備", "2", "容逸燊、陳昆仁、祥維、囿羽"],
        ["被動元件／MLCC", "1", "郭哲榮、林漢偉、陳昆仁、鍾國忠、黃豐凱（08/28 才成形）"],
        ["矽晶圓", "1", "林漢偉、陳昆仁、郭哲榮、楊雲翔、祥維（08/24，之後轉為分歧）"],
    ]
    widths = [Inches(2.7), Inches(1.9), Inches(7.57)]
    tbl = table(slide, MARGIN, Inches(1.42), BODY_W, widths, rows, row_height=Inches(0.42))
    for r in range(1, 7):
        for run in tbl.cell(r, 0).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
        for run in tbl.cell(r, 1).text_frame.paragraphs[0].runs:
            run.font.name = MONO
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = UP if r <= 3 else MUTED
    # 天數長條
    for r, days in ((1, 4), (2, 4), (3, 4), (4, 2), (5, 1), (6, 1)):
        top = Inches(1.42 + 0.32 + 0.42 * (r - 1) + 0.15)
        rect(slide, MARGIN + Inches(3.05), top, Inches(0.24 * days),
             Pt(8), fill=UP if days >= 4 else UP_SOFT)

    card(slide, MARGIN, Inches(4.55), BODY_W, Inches(1.75))
    frame = textbox(slide, MARGIN + Inches(0.28), Inches(4.75), BODY_W - Inches(0.56), Inches(0.3))
    write(frame, "跨日重複出現的個股（括號為週內被列入一致看多的天數）", size=15, color=NAVY,
          bold=True, first=True, space_after=0)
    chips = [("大立光", "3008", 4), ("聯亞", "3081", 3), ("南亞科", "2408", 3),
             ("南電", "8046", 2), ("景碩", "3189", 2), ("聯茂", "6213", 2),
             ("聯電", "2303", 2), ("光寶科", "2301", 2), ("環球晶", "6488", 2),
             ("華星光", "4979", 2)]
    for i, (name, code, days) in enumerate(chips):
        col, row = i % 5, i // 5
        left = MARGIN + Inches(0.28) + Inches(col * 2.42)
        top = Inches(5.25 + row * 0.52)
        rect(slide, left, top, Inches(2.24), Inches(0.42), fill=WHITE,
             line=UP if days >= 4 else (UP_SOFT if days >= 3 else LINE))
        frame = textbox(slide, left + Inches(0.12), top + Inches(0.08), Inches(2.0), Inches(0.3))
        para = frame.paragraphs[0]
        para.space_after = Pt(0)
        for text, size, color, bold, mono in (
            (name + "　", 13, INK, days >= 3, False),
            (code + "　", 10.5, MUTED, False, True),
            (str(days), 13, UP if days >= 3 else MUTED, True, True),
        ):
            run = para.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = MONO if mono else SANS
    footer(slide, "天數為五個交易日內被列入「一致看多」章節的次數，非漲跌統計。")


def consensus_down(prs):
    slide = blank(prs)
    band(slide, "04", "看空與風險標的", "週內唯一方向反轉的族群，在這一頁")

    rect(slide, MARGIN, Inches(1.45), Pt(3), Inches(0.3), fill=DOWN)
    frame = textbox(slide, MARGIN + Inches(0.16), Inches(1.42), Inches(6.1), Inches(0.4))
    write(frame, "貨櫃航運：一週內從共識看多翻成共識看空", size=19, color=INK, bold=True,
          first=True, space_after=0)

    steps = [
        ("08/24–25", UP, "看多", "林漢偉、張林忠、老王 —— 塞港短期無解、低本益比高殖利率",
         "郭哲榮同期一路看空，認為與 AI 無關"),
        ("08/26", DOWN, "翻空", "萊茵河水位回升、歐線期貨大跌，林漢偉當日「全數出清」，"
                                "與郭哲榮、容逸燊形成一致看空",
         "同日萬海（2615）、長榮（2603）被列入一致看空個股"),
        ("08/27", MUTED, "中性", "陳昆仁僅視為非 AI 的強勢概念股，低進高出", None),
    ]
    top = Inches(1.95)
    for label, color, verdict, body, note in steps:
        height = Inches(1.18) if note else Inches(0.86)
        card(slide, MARGIN, top, Inches(6.1), height, accent=color, side="left")
        frame = textbox(slide, MARGIN + Inches(0.24), top + Inches(0.16), Inches(1.1), Inches(0.3))
        write(frame, label, size=13, color=color, bold=True, mono=True, first=True, space_after=0)
        frame = textbox(slide, MARGIN + Inches(1.4), top + Inches(0.14), Inches(4.5),
                        height - Inches(0.28))
        para = frame.paragraphs[0]
        para.space_after = Pt(4)
        run = para.add_run()
        run.text = verdict + "　"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = SANS
        run = para.add_run()
        run.text = body
        run.font.size = Pt(12.5)
        run.font.color.rgb = INK
        run.font.name = SANS
        if note:
            write(frame, note, size=11.5, color=MUTED, space_after=0)
        top = top + height + Inches(0.12)

    frame = textbox(slide, MARGIN, Inches(5.5), Inches(6.1), Inches(0.6))
    write(frame, "同一族群、同一位分析師，兩天內方向相反 —— 這是逐日摘要才看得出來的東西。",
          size=12.5, color=MUTED, first=True, space_after=0)

    right = Inches(7.1)
    right_w = SLIDE_W - right - MARGIN
    rect(slide, right, Inches(1.45), Pt(3), Inches(0.3), fill=UP)
    frame = textbox(slide, right + Inches(0.16), Inches(1.42), right_w, Inches(0.4))
    write(frame, "欣興（3037）：週五爆出的個股事件", size=19, color=INK, bold=True,
          first=True, space_after=0)

    card(slide, right, Inches(1.95), right_w, Inches(2.05))
    frame = textbox(slide, right + Inches(0.24), Inches(2.14), right_w - Inches(0.48), Inches(1.7))
    write(frame, "兩名副總與一名會計遭檢調帶走，搜索規模大，引發內控與治理爭議。",
          size=13, color=INK, first=True, space_after=7)
    write(frame, "看空者：陳昆仁、黃豐凱、郭哲榮", size=14, color=DOWN, bold=True, space_after=7)
    write(frame, "黃豐凱：投信高持股且週五來不及調節，預估跌至 900–1,000 元；"
                 "資金可能轉單至景碩（3189）避險", size=12.5, color=INK, space_after=5)
    write(frame, "郭哲榮：若僅個人內線是一天反應；若涉及假帳、營收灌水，可能連續跌停兩至三根",
          size=12.5, color=INK, space_after=0)

    frame = textbox(slide, right, Inches(4.2), right_w, Inches(0.3))
    write(frame, "其他看空項目", size=15, color=NAVY, bold=True, first=True, space_after=0)
    others = [
        ("消費性電子", "漲價壓力侵蝕購買力、拉貨不積極（容逸燊、游庭皓、蔚辰）", None),
        ("成熟製程晶圓代工", "老王、郭哲榮、陳昆仁、祥維看空；容逸燊、黃豐凱在 08/28 看多聯電",
         "全週未收斂"),
    ]
    top = Inches(4.62)
    for title, body, tag in others:
        height = Inches(0.62) if tag is None else Inches(0.86)
        card(slide, right, top, right_w, height)
        frame = textbox(slide, right + Inches(0.22), top + Inches(0.12), right_w - Inches(0.44),
                        height - Inches(0.2))
        para = frame.paragraphs[0]
        para.space_after = Pt(3)
        run = para.add_run()
        run.text = title + "　"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = INK
        run.font.name = SANS
        run = para.add_run()
        run.text = body
        run.font.size = Pt(12)
        run.font.color.rgb = INK
        run.font.name = SANS
        if tag:
            write(frame, tag, size=12, color=UP, bold=True, space_after=0)
        top = top + height + Inches(0.12)
    footer(slide, "個股事件為節目中的轉述，法律程序與後續發展未在來源中確認。")


def disagreements(prs):
    slide = blank(prs)
    band(slide, "05", "貫穿整週的分歧", "分歧比共識更值得記，因為下週會由其中一方勝出")

    heads = [("議題", MUTED, MARGIN + Inches(0.2), Inches(2.1)),
             ("一方", UP, MARGIN + Inches(2.5), Inches(4.9)),
             ("另一方", DOWN, MARGIN + Inches(7.6), Inches(4.5))]
    for text, color, left, width in heads:
        frame = textbox(slide, left, Inches(1.4), width, Inches(0.3))
        write(frame, text, size=12, color=color, bold=True, first=True, space_after=0)
    rect(slide, MARGIN, Inches(1.72), BODY_W, Pt(2), fill=NAVY)

    items = [
        ("矽晶圓基本面", None,
         "郭哲榮、劉育綸：現貨雙位數大漲、產能滿載、合約洽談中，17 天內信號全面轉正",
         "蔡明翰：營收與獲利在所有族群中最弱，純屬預期，拉回「絕對不要向下攤平」"),
        ("高檔熱門股要不要追", None,
         "林漢偉：順著強勢產業趨勢做，領頭羊先創高，背後有大資金帶動",
         "容逸燊：聯亞、大立光短線已拉數根漲停，賺賠比不划算，轉向低位階個股"),
        ("國巨（2327）", "被動元件",
         "郭哲榮：1,000 元以上有 30 萬人套牢，反覆洗盤，拉回即買點",
         "老王、劉育綸（08/25）：跌破三短均線「三聲無奈」，土洋同賣、散戶攤平為反指標"),
        ("記憶體的政治風險", None,
         "林漢偉、高憲容：輝達親口證實記憶體極度缺貨，外資單月敲進逾 10 萬張",
         "陳昆仁：9/24 川習會若美方對中國記憶體鬆綁，南亞科、華邦電、旺宏恐重挫"),
    ]
    top = Inches(1.9)
    height = Inches(0.94)
    for title, sub, left_text, right_text in items:
        card(slide, MARGIN, top, BODY_W, height)
        frame = textbox(slide, MARGIN + Inches(0.24), top + Inches(0.2), Inches(2.1),
                        height - Inches(0.3))
        write(frame, title, size=14.5, color=INK, bold=True, first=True, space_after=2)
        if sub:
            write(frame, sub, size=11.5, color=MUTED, space_after=0)
        rect(slide, MARGIN + Inches(2.44), top, Pt(3), height, fill=UP)
        frame = textbox(slide, MARGIN + Inches(2.66), top + Inches(0.18), Inches(4.7),
                        height - Inches(0.3))
        write(frame, left_text, size=12.5, color=INK, first=True, space_after=0)
        rect(slide, MARGIN + Inches(7.5), top, Pt(3), height, fill=DOWN)
        frame = textbox(slide, MARGIN + Inches(7.72), top + Inches(0.18), Inches(4.35),
                        height - Inches(0.3))
        write(frame, right_text, size=12.5, color=INK, first=True, space_after=0)
        top = top + height + Inches(0.12)

    rect(slide, MARGIN, Inches(6.35), Inches(0.14), Inches(0.14), fill=NAVY)
    frame = textbox(slide, MARGIN + Inches(0.3), Inches(6.28), Inches(11.8), Inches(0.34))
    write(frame, "被動元件是唯一在週內收斂的分歧：08/25 還在吵，08/28 變成當日看多人數最多的族群。",
          size=13.5, color=INK, first=True, space_after=0)
    footer(slide, "分歧雙方的敘述均取自各自節目，未做對錯判定。")


def calendar(prs):
    slide = blank(prs)
    band(slide, "06", "下週起的關鍵時程", "節目中被點名的日期，依時間排序")
    items = [
        ("08/31 (一)", UP, "MSCI 季度調整生效", "林漢偉：補量關鍵日，預估成交量 1.1 至 1.2 兆元"),
        ("08/31 (一)", UP, "矽光子國際論壇", "陳昆仁：CPO 商轉元年，提前布局"),
        ("09/02–04", NAVY, "台北國際半導體展", "陳昆仁、容逸燊、陳唯泰：設備與 CPO 族群題材期"),
        ("09/03 (四)", NAVY, "中光電投控（3718）換股掛牌", "鍾國忠"),
        ("09/09 (三)", NAVY, "美國財政部加倍美債回購",
         "游庭皓、林漢偉：等同變相寬鬆，利於台幣升值與外資回流"),
        ("09/10 (四)", NAVY, "蘋果 iPhone 18 系列與折疊機發表會",
         "林漢偉、陳昆仁、陳唯泰：光學、折疊機鉸鏈、軸承族群"),
        ("09/15 (二)", NAVY, "台北航太展", "鍾國忠：無人機預算商機"),
        ("09/24 (四)", DOWN, "川習會", "游庭皓、陳昆仁：記憶體與關稅的政治風險點"),
        ("九月內", LINE, "投信季底作帳、Fed 會議、日銀升息", "王建文、紀緯明、鍾國忠"),
    ]
    top = Inches(1.42)
    height = Inches(0.53)
    for i, (date, color, title, who) in enumerate(items):
        card(slide, MARGIN, top, BODY_W, height,
             accent=color if color is not LINE else LINE, side="left")
        frame = textbox(slide, MARGIN + Inches(0.22), top + Inches(0.14), Inches(1.5), Inches(0.3))
        write(frame, date, size=13.5, color=color if color is not LINE else MUTED, bold=True,
              mono=(color is not LINE), first=True, space_after=0)
        frame = textbox(slide, MARGIN + Inches(1.85), top + Inches(0.13), Inches(3.6), Inches(0.32))
        write(frame, title, size=14, color=INK, bold=True, first=True, space_after=0)
        frame = textbox(slide, MARGIN + Inches(5.6), top + Inches(0.15), Inches(6.4), Inches(0.32))
        write(frame, who, size=12.5, color=INK, first=True, space_after=0)
        top = top + height + Inches(0.075)
    footer(slide, "日期以節目中所述為準；跨日出現差異者取最後一次提及的版本。")


def next_week(prs):
    slide = blank(prs)
    band(slide, "07", "下週判斷：多空兩邊各自的條件", "沒有結論，只有各自畫出來的線")

    col_w = Inches(6.05)
    for i, (color, title, entries) in enumerate((
        (UP, "偏多方", [
            ("林漢偉", "看多／信心高",
             "站穩 46,325 點，均線重回多頭排列，九月第一週有機會突破 48,218 點"),
            ("高憲容", "看多／信心高",
             "台幣轉升 ＋ 外資史上最大買超，依 2022/11、2023/11、2025/5 三次經驗，"
             "可啟動 7–10 個月、漲幅 38%–105% 的波段，台股看 54,000 點"),
            ("郭哲榮", "看多／信心高", "下週被動元件大補漲；年底看 5 萬點"),
            ("容逸燊", "看多／信心高", "九月資金轉向低基期的工具機、被動元件、成熟製程"),
        ]),
        (DOWN, "偏空與保留方", [
            ("鍾國忠", "偏空／信心中",
             "46,350 點以上追高將面臨巨大賣壓，需 1.1 兆換手量；建議逢高減碼、保留現金"),
            ("囿羽", "偏多／信心中",
             "量能與 KD 雙背離、季線仍下彎、下週布林通道高檔收窄，觀察日 KD 能否化解背離"),
            ("劉育綸", "中性／信心中",
             "若跌破當日低點，將進入 44,200–46,400 點的箱型整理，「創高不追、回落敢買」"),
            ("陳昆仁", "偏多／信心中", "前波七月套牢量大，短線有震盪與獲利了結賣壓"),
        ]),
    )):
        left = MARGIN + Inches(i * 6.28)
        rect(slide, left, Inches(1.42), col_w, Inches(0.42), fill=color)
        frame = textbox(slide, left + Inches(0.2), Inches(1.5), col_w - Inches(0.4), Inches(0.3))
        write(frame, title, size=16, color=WHITE, bold=True, first=True, space_after=0)
        top = Inches(1.92)
        for name, stance, body in entries:
            height = Inches(0.92) if len(body) > 40 else Inches(0.72)
            card(slide, left, top, col_w, height)
            frame = textbox(slide, left + Inches(0.22), top + Inches(0.12), col_w - Inches(0.44),
                            height - Inches(0.2))
            para = frame.paragraphs[0]
            para.space_after = Pt(4)
            run = para.add_run()
            run.text = name + "　"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = INK
            run.font.name = SANS
            run = para.add_run()
            run.text = stance
            run.font.size = Pt(11.5)
            run.font.color.rgb = color
            run.font.name = SANS
            write(frame, body, size=12.5, color=INK, space_after=0)
            top = top + height + Inches(0.08)

    rect(slide, MARGIN, Inches(5.72), BODY_W, Inches(1.16), fill=NAVY)
    frame = textbox(slide, MARGIN + Inches(0.28), Inches(5.86), Inches(4.4), Inches(0.34))
    write(frame, "多空雙方其實在講同一件事：量。", size=16, color=WHITE, bold=True, first=True,
          space_after=0)
    volumes = [("8,500 億–1 兆", "劉育綸　進攻訊號"), ("1.1 兆", "鍾國忠　換手門檻"),
               ("1.1–1.2 兆", "林漢偉　MSCI 補量"), ("1.2 兆", "容逸燊　站穩新高")]
    for i, (value, who) in enumerate(volumes):
        left = MARGIN + Inches(4.9) + Inches(i * 1.85)
        rect(slide, left, Inches(5.86), Pt(2), Inches(0.86), fill=UP)
        frame = textbox(slide, left + Inches(0.14), Inches(5.86), Inches(1.66), Inches(0.86))
        write(frame, value, size=15, color=WHITE, bold=True, mono=True, first=True, space_after=3)
        write(frame, who, size=11, color=NAVY_SOFT, space_after=0)
    footer(slide, "態度與信心（看多／偏多／中性／偏空／看空、高／中／低）為來源摘要中的原始標註。")


def sources(prs):
    slide = blank(prs)
    band(slide, "08", "資料來源與整理原則")

    left_w = Inches(6.05)
    frame = textbox(slide, MARGIN, Inches(1.42), left_w, Inches(0.3))
    write(frame, "來源檔案", size=16, color=NAVY, bold=True, first=True, space_after=0)
    card(slide, MARGIN, Inches(1.82), left_w, Inches(1.35))
    frame = textbox(slide, MARGIN + Inches(0.24), Inches(1.98), left_w - Inches(0.48), Inches(1.1))
    write(frame, "Stock_Summary/20260824_StockSummary.txt", size=12, color=INK, mono=True,
          first=True, space_after=3)
    write(frame, "…　20260825　20260826　20260827　…", size=12, color=MUTED, mono=True,
          space_after=3)
    write(frame, "Stock_Summary/20260828_StockSummary.txt", size=12, color=INK, mono=True,
          space_after=6)
    write(frame, "每份為當日多支財經節目的彙整摘要，五份合計約 12 萬字。", size=11.5, color=MUTED,
          space_after=0)

    frame = textbox(slide, MARGIN, Inches(3.35), left_w, Inches(0.3))
    write(frame, "整理原則（沿用同一份摘要規則書）", size=16, color=NAVY, bold=True, first=True,
          space_after=0)
    rules = [
        "只寫來源明確講過的內容，不推論、不補充來源之外的知識",
        "每一條都標明是哪位分析師講的；同一人在多個頻道發言合併計算一次",
        "態度只用看多／偏多／中性／偏空／看空五個詞，信心分高／中／低",
        "來源未查證的股名或代號一律保留「待確認」，不憑印象補上",
    ]
    top = Inches(3.75)
    for rule in rules:
        card(slide, MARGIN, top, left_w, Inches(0.46))
        frame = textbox(slide, MARGIN + Inches(0.22), top + Inches(0.12), left_w - Inches(0.44),
                        Inches(0.3))
        write(frame, rule, size=12.5, color=INK, first=True, space_after=0)
        top = top + Inches(0.55)

    right = MARGIN + Inches(6.28)
    right_w = SLIDE_W - right - MARGIN
    frame = textbox(slide, right, Inches(1.42), right_w, Inches(0.3))
    write(frame, "已知的來源限制", size=16, color=UP, bold=True, first=True, space_after=0)
    limits = [
        "08/28 的來源檔在最後一位分析師的個股段落中斷（新盛力 3211 之後無內容），該段未納入本簡報",
        "來源檔在 08/25 寫「蔡明漢」、08/26 寫「蔡明翰」，經確認為同一人，正確為蔡明翰，本簡報已統一",
        "指數、成交量、財務數字均為節目中所述，未另行向公開資訊源查證",
    ]
    top = Inches(1.82)
    for text in limits:
        card(slide, right, top, right_w, Inches(0.72), accent=UP, side="left")
        frame = textbox(slide, right + Inches(0.24), top + Inches(0.14), right_w - Inches(0.48),
                        Inches(0.5))
        write(frame, text, size=12.5, color=INK, first=True, space_after=0)
        top = top + Inches(0.82)

    rect(slide, right, Inches(4.5), right_w, Inches(1.75), fill=NAVY)
    rect(slide, right + Inches(0.26), Inches(4.76), Inches(0.14), Inches(0.14), fill=UP)
    frame = textbox(slide, right + Inches(0.56), Inches(4.68), right_w - Inches(0.8), Inches(1.3))
    write(frame, "本簡報為節目內容整理", size=17, color=WHITE, bold=True, first=True, space_after=6)
    write(frame, "不構成任何投資建議。", size=16, color=PINK, bold=True, space_after=8)
    write(frame, "內容為講者說法的轉述，非整理者的看法。投資決策請自行查證並承擔風險。",
          size=11.5, color=NAVY_SOFT, space_after=0)

    footer(slide, "2026/08/24 – 08/28　台股財經節目週摘要　·　Generated with Claude Code")


def main() -> None:
    prs = new_deck()
    cover(prs)
    week_track(prs)
    macro(prs)
    consensus_up(prs)
    consensus_down(prs)
    disagreements(prs)
    calendar(prs)
    next_week(prs)
    sources(prs)
    prs.save(OUT_FILE)
    print(f"[ppt] {OUT_FILE.name}：{len(prs.slides._sldIdLst)} 頁")


if __name__ == "__main__":
    main()
